"""
Sentinel — investor / AI Kavach competition deck builder.

Two stages:
  1) Pillow renders premium dark "cyber" slide backgrounds (ink gradient,
     faint hex grid, soft cyan/magenta corner glows) and a cyan->magenta
     rule used as a title underline. Baked to assets/brand/.
  2) python-pptx assembles a ~29-slide 16:9 deck: title, section dividers,
     content bullets, framed REAL CLI screenshots (assets/brand/shot_*.png,
     produced by capture_cli.py + render_shots.py against the live engine),
     a validation stat band, and a category-level market matrix. Real speaker
     notes on every slide.

Run with a Python that has Pillow + python-pptx installed:
    ./.venv_docs/bin/python assets/build_deck.py
Output: deck/Sentinel_AI_Kavach.pptx
"""

import math
import os

from PIL import Image, ImageDraw, ImageFilter

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE = os.path.dirname(__file__)
ROOT = os.path.dirname(HERE)
BRAND = os.path.join(HERE, "brand")
DECK = os.path.join(ROOT, "deck")
os.makedirs(DECK, exist_ok=True)

# ---- palette ----
CYAN = (34, 211, 238)
MAGENTA = (232, 75, 255)
INK_TOP = (11, 16, 26)
INK_BOT = (6, 9, 15)
GRID = (18, 27, 42)
HEXDIM = (24, 34, 52)

# pptx colors
C_CYAN = RGBColor(0x22, 0xD3, 0xEE)
C_MAGENTA = RGBColor(0xE8, 0x4B, 0xFF)
C_WHITE = RGBColor(0xEA, 0xF2, 0xFF)
C_BODY = RGBColor(0xC4, 0xD2, 0xE4)
C_DIM = RGBColor(0x6B, 0x7A, 0x90)
C_GREEN = RGBColor(0x2D, 0xD4, 0x8F)
C_AMBER = RGBColor(0xF5, 0xB3, 0x01)
C_RED = RGBColor(0xFF, 0x5C, 0x6E)

W, H = 1920, 1080


# ---------- background rendering ----------
def _vgrad(size, c0, c1):
    w, h = size
    base = Image.new("RGB", (1, h))
    px = base.load()
    for y in range(h):
        t = y / max(1, h - 1)
        px[0, y] = tuple(int(c0[i] + (c1[i] - c0[i]) * t) for i in range(3))
    return base.resize((w, h))


def _hexagon(cx, cy, r, rot=-math.pi / 2):
    return [(cx + r * math.cos(rot + k * math.pi / 3),
             cy + r * math.sin(rot + k * math.pi / 3)) for k in range(6)]


def _glow(center, radius, color, alpha):
    layer = Image.new("RGBA", (W, H), (0, 0, 0, 0))
    d = ImageDraw.Draw(layer)
    cx, cy = center
    d.ellipse([cx - radius, cy - radius, cx + radius, cy + radius],
              fill=color + (alpha,))
    return layer.filter(ImageFilter.GaussianBlur(radius * 0.55))


def _base_canvas():
    base = _vgrad((W, H), INK_TOP, INK_BOT).convert("RGBA")
    grid = Image.new("RGBA", (W, H), (0, 0, 0, 0))
    gd = ImageDraw.Draw(grid)
    step = 54
    for x in range(0, W, step):
        gd.line([(x, 0), (x, H)], fill=GRID + (255,), width=1)
    for y in range(0, H, step):
        gd.line([(0, y), (W, y)], fill=GRID + (255,), width=1)
    base = Image.alpha_composite(base, grid)
    hx = Image.new("RGBA", (W, H), (0, 0, 0, 0))
    hd = ImageDraw.Draw(hx)
    for cx, cy, r in ((-60, H + 80, 460), (W + 90, -70, 420)):
        poly = _hexagon(cx, cy, r)
        hd.line(poly + [poly[0]], fill=HEXDIM + (255,), width=3, joint="curve")
    return Image.alpha_composite(base, hx)


def render_content_bg(path):
    base = _base_canvas()
    base = Image.alpha_composite(base, _glow((150, 120), 520, CYAN, 34))
    base = Image.alpha_composite(base, _glow((W - 120, H - 90), 560, MAGENTA, 30))
    spine = Image.new("RGBA", (W, H), (0, 0, 0, 0))
    sd = ImageDraw.Draw(spine)
    sd.rectangle([0, 0, 8, H], fill=CYAN + (150,))
    base = Image.alpha_composite(base, spine)
    base.convert("RGB").save(path, quality=95)


def render_title_bg(path):
    base = _base_canvas()
    base = Image.alpha_composite(base, _glow((W // 2 - 260, H // 2 - 40), 620, CYAN, 40))
    base = Image.alpha_composite(base, _glow((W // 2 + 300, H // 2 + 60), 640, MAGENTA, 38))
    scan = Image.new("RGBA", (W, H), (0, 0, 0, 0))
    scd = ImageDraw.Draw(scan)
    scd.rectangle([0, H // 2 - 2, W, H // 2 + 2], fill=CYAN + (26,))
    base = Image.alpha_composite(base, scan.filter(ImageFilter.GaussianBlur(2)))
    base.convert("RGB").save(path, quality=95)


def render_divider_bg(path):
    base = _base_canvas()
    base = Image.alpha_composite(base, _glow((W // 2, H // 2), 720, MAGENTA, 34))
    base = Image.alpha_composite(base, _glow((W // 2 - 380, H // 2), 520, CYAN, 26))
    base.convert("RGB").save(path, quality=95)


def render_rule(path):
    w, h = 1200, 22
    bar = Image.new("RGB", (w, 1))
    px = bar.load()
    for x in range(w):
        t = x / (w - 1)
        px[x, 0] = tuple(int(CYAN[i] + (MAGENTA[i] - CYAN[i]) * t) for i in range(3))
    bar = bar.resize((w, h)).convert("RGBA")
    mask = Image.new("L", (w, h), 0)
    ImageDraw.Draw(mask).rounded_rectangle([0, 0, w - 1, h - 1], radius=h // 2, fill=255)
    bar.putalpha(mask)
    bar.save(path)
# ---------- slide content ----------
# schema: title, subtitle, notes; optional kind ∈ {title,divider,content};
#   bullets (str | ("ok"/"next"/"cyan"/"warn", text)); image (+caption,layout);
#   stats [(value,label)]; cols + rows [(label,[yes/partial/no,...])].
SLIDES = []
SLIDES += [
    {
        "kind": "title",
        "title": "SENTINEL",
        "subtitle": "Autonomous cyber-reasoning  —  find → reason → prove → patch → prove",
        "bullets": [
            "Local-first AI agent that reasons about live web targets",
            "A status code is never a verdict — a deterministic judge decides",
            "Proven end-to-end on two independent live stacks",
            "Ten vulnerability classes, each closing the full loop live",
            "Every decision explainable, auditable, and scope-guarded",
        ],
        "notes": (
            "Sentinel is an autonomous cyber-reasoning agent. Point it at a live target; "
            "it recons, hypothesizes, probes under a strict scope guard, and lets a "
            "deterministic judge decide every verdict — then it patches and re-proves the fix. "
            "Local-first, evidence-driven, target-agnostic."
        ),
    },
    {
        "title": "Scanners Guess. Sentinel Reasons.",
        "subtitle": "From noisy 200s to justified, reproduced, remediated claims",
        "bullets": [
            "Signature scanners flag 200s — a firehose of false positives",
            "Authorization & session flaws need context, not patterns",
            "Teams drown in unranked output nobody can trust or triage",
            "A finding should mean a reproduced contradiction — nothing less",
            "And the fix should be proven, not just recommended",
        ],
        "notes": (
            "The core pain: scanners pattern-match and treat any 200 as a hit, burying teams "
            "in false positives — and the highest-impact classes (authorization, session) are "
            "exactly what signatures miss. Sentinel's bar is a reproduced contradiction plus a "
            "proven fix."
        ),
    },
]
SLIDES += [
    {
        "title": "Why Now",
        "subtitle": "The attack surface exploded; verification didn't keep up",
        "bullets": [
            "APIs and microservices multiplied the authorization surface",
            "Broken access control is the #1 web risk (OWASP Top 10)",
            "Capable local LLMs make on-prem autonomous reasoning viable",
            "But raw LLM output is unverifiable — and legally radioactive",
            "The unlock: bounded AI reasoning + a deterministic ground truth",
        ],
        "notes": (
            "Two curves cross: the authorization surface exploded with APIs while local models "
            "got good enough to reason on-prem. The missing piece is trust — an LLM alone can't "
            "be the authority. Sentinel supplies the deterministic ground truth that makes "
            "autonomous reasoning safe to ship."
        ),
    },
    {
        "title": "The Insight — Reason, Don't Signature",
        "subtitle": "Separate discovery from judgement; cage the AI",
        "bullets": [
            ("cyan", "Discovery produces hypotheses — explicitly NOT findings"),
            ("cyan", "A pure, deterministic judge owns every verdict"),
            ("cyan", "The LLM is advisory only: it breaks ties, never decides"),
            ("cyan", "Findings crystallize only from reproduced contradictions"),
            ("cyan", "Fixes are proven live — the same judge must flip its verdict"),
        ],
        "notes": (
            "Our thesis in one line: reason, don't signature. Discovery, judgement, and "
            "remediation are separate stages. The AI accelerates search but is caged out of the "
            "verdict; only the deterministic judge confirms — and only a live re-proof closes a "
            "finding."
        ),
    },
    {
        "title": "What Sentinel Is",
        "subtitle": "recon → hypotheses → probe → judge → patch → re-prove",
        "bullets": [
            "Recon materializes a security graph of the live surface",
            "Hypotheses = justified, conservative reasons to test",
            "Scope-guarded probes record facts, never verdicts",
            "The deterministic judge tests facts against an operator oracle",
            "Confirmed findings are patched on a shield and re-proven live",
        ],
        "notes": (
            "End to end it is a disciplined loop: graph the surface, seed hypotheses, fire one "
            "scope-checked probe at a time, judge the observations against explicit policy, then "
            "synthesize and enforce a fix and re-run the same judge to prove it holds."
        ),
    },
]
SLIDES += [
    {"kind": "divider", "title": "The Epistemic Contract",
     "subtitle": "Why a Sentinel finding is worth believing",
     "notes": "The heart of the product: a set of hard rules that make every verdict "
              "trustworthy and every fix real."},
    {
        "title": "The Epistemic Contract",
        "subtitle": "Hard rules the engine cannot violate",
        "bullets": [
            ("ok", "A finding requires the PURE judge to reproduce a contradiction"),
            ("ok", "vs an operator-declared oracle — never a status code alone"),
            ("ok", "FIX PROVEN requires that SAME judge to flip under live enforcement"),
            ("warn", "The AI never manufactures a finding, verdict, or evidence"),
            ("warn", "A compliant control → DISPROVED → no finding (honest differential)"),
        ],
        "notes": (
            "These invariants are enforced in code. Verification runs the pure judge on a scratch "
            "graph seeded with relationships only — never the confirmation path — so a fix is "
            "proven by a real VALIDATED→DISPROVED flip. Compliant controls yield no finding, "
            "which is how you know it isn't fabricating."
        ),
    },
    {
        "title": "The Deterministic Judge",
        "subtitle": "The caged advisor — smart search, bounded authority",
        "bullets": [
            "A transparent score ranks candidates; information gain decays",
            "The local LLM is consulted only to break top ties",
            "It picks by list index mapped to a real id — can't invent targets",
            "Every judgement ships a human-readable rationale trail",
            "AI reasoning is shown as telemetry, never as authority",
        ],
        "notes": (
            "Our answer to 'isn't it just an LLM wrapper?' The engine scores candidates with a "
            "self-adapting formula and consults the model only on ties, where it chooses by index "
            "we map back to a real id. It cannot hallucinate a target or override a higher score."
        ),
    },
    {
        "title": "Architecture",
        "subtitle": "A security graph, a scope guard, and a loopback shield",
        "bullets": [
            "Security graph: hosts, resources, principals, evidence, findings",
            "Pre-connection scope guard refuses out-of-scope probes",
            "Ten pure judges: authz · posture · cookies · privesc · SQLi · SSTI · redirect · CORS · SSRF · JWT",
            "RemediationEnforcer: an ephemeral loopback reverse-proxy shield",
            "Local model via Ollama (qwen3:4b) — no cloud, no data egress",
        ],
        "notes": (
            "The graph is the single source of truth; the scope guard is checked before any "
            "socket opens; each vuln class has its own pure judge; and the enforcer stands up a "
            "loopback shield to prove fixes. All local — the model runs on-box."
        ),
    },
]
SLIDES += [
    {"kind": "divider", "title": "Live, End-to-End",
     "subtitle": "Real screenshots — every verdict produced by the engine",
     "notes": "The following screenshots are captured from the real engine driving a live "
              "target over real sockets. Nothing is staged or mocked."},
    {
        "title": "FIND — Recon → Oracle → Hypotheses",
        "subtitle": "It models the surface, then reasons about it",
        "image": "shot_recon.png", "layout": "full",
        "caption": "Live capture: recon graph, the operator oracle, and seeded hypotheses.",
        "notes": (
            "The opening move: Sentinel recons the target into a security graph, loads the "
            "operator's access-policy oracle, and seeds conservative hypotheses. Discovery is "
            "explicitly not a vulnerability — each hypothesis is a justified reason to test."
        ),
    },
    {
        "title": "The Decision Board",
        "subtitle": "Autonomous research cycles, fully auditable",
        "image": "shot_cycles.png", "layout": "right",
        "bullets": [
            "One scope-checked probe per cycle",
            "Deterministic score picks the next test",
            "Advisor consulted only on ties",
            "Every choice carries a rationale",
            "Dead leads decay and are abandoned",
        ],
        "notes": (
            "The decision board shows the loop thinking out loud: candidates, scores, the "
            "occasional advisory tie-break, and the probe fired. It reads like a researcher's "
            "notebook, not a black box."
        ),
    },
    {
        "title": "CONFIRMED — Broken Access Control",
        "subtitle": "A reproduced contradiction, not a guess",
        "image": "shot_findings.png", "layout": "right",
        "bullets": [
            ("cyan", "Oracle: anonymous MUST be denied /api/Feedbacks"),
            ("cyan", "Observed: the whole collection leaks (200)"),
            ("cyan", "Judge reproduces the contradiction → CONFIRMED"),
            ("cyan", "The compliant control (/api/Users, 401) → no finding"),
            ("cyan", "Confidence and provenance travel with the finding"),
        ],
        "notes": (
            "A HIGH-severity broken-access-control finding: the oracle says deny, the live target "
            "leaks, and the judge reproduces it. The neighbouring compliant control produces no "
            "finding — the honest differential that proves nothing is manufactured."
        ),
    },
    {
        "title": "PATCH + PROVE",
        "subtitle": "The fix is proven live, not recommended",
        "image": "shot_patch_prove.png", "layout": "right",
        "bullets": [
            ("ok", "Synthesize a corrective control for the finding"),
            ("ok", "Enforce it on an ephemeral loopback shield"),
            ("ok", "Re-run the SAME pure judge through the shield"),
            ("ok", "before 200 VALIDATED → after 403 DISPROVED = FIX PROVEN"),
            ("ok", "Portable artifacts: nginx · envoy · caddy · JSON"),
        ],
        "notes": (
            "Remediation is a proof, not a suggestion. Sentinel enforces the control on a "
            "loopback shield and re-runs the identical judge; only a real VALIDATED→DISPROVED "
            "flip earns FIX PROVEN, and it emits config artifacts for real proxies."
        ),
    },
]
SLIDES += [
    {"kind": "divider", "title": "Ten Vulnerability Classes",
     "subtitle": "Each closes the full find → patch → prove loop, live",
     "notes": "Sentinel is not a one-trick demo. Ten independent classes each run the entire "
              "loop end-to-end against live targets — every one gated by the same epistemic contract."},
    {
        "title": "Class 1 — Broken Access Control",
        "subtitle": "authorization_policy_violation",
        "bullets": [
            "Operator declares who may do what (the access oracle)",
            "Sentinel probes each rule as the declared principal",
            "Judge confirms only reproduced authorization contradictions",
            "Authenticated mode re-points rules to a captured session",
            "Fix: deny/allow enforcement proven on the shield",
        ],
        "notes": (
            "The flagship class. The oracle encodes intended authorization; Sentinel tests it "
            "live and confirms only reproduced violations. With the Login Tester it runs the "
            "same rules as an authenticated principal."
        ),
    },
    {
        "title": "Class 2 — Security-Header Posture",
        "subtitle": "security_misconfiguration",
        "image": "shot_posture.png", "layout": "right",
        "bullets": [
            ("cyan", "Declared posture: CSP, CORS, Referrer-Policy, …"),
            ("cyan", "Judge reads the LIVE response headers"),
            ("cyan", "Absent CSP, wildcard CORS → CONFIRMED"),
            ("cyan", "Present nosniff control → DISPROVED, no finding"),
            ("ok", "Shield injects/strips headers → re-proven"),
        ],
        "notes": (
            "Header posture as a first-class, judged property. Missing CSP and wildcard CORS are "
            "confirmed against the live response; a compliant X-Content-Type-Options yields no "
            "finding; then the shield fixes it and the judge re-proves."
        ),
    },
    {
        "title": "Class 3 — Insecure Cookies",
        "subtitle": "insecure_cookie — the session-theft pivot",
        "image": "shot_cookie.png", "layout": "right",
        "bullets": [
            ("cyan", "Grounded in OBSERVED Set-Cookie — never guessed"),
            ("cyan", "Missing HttpOnly / Secure, weak SameSite"),
            ("cyan", "The classic XSS→session-theft chaining ingredient"),
            ("ok", "Shield rewrites Set-Cookie → flags proven present"),
            ("ok", "Artifacts for nginx proxy_cookie_flags, envoy, caddy"),
        ],
        "notes": (
            "The newest class and a chaining enabler. Cookie oracles are grounded in the Set-"
            "Cookie the target actually sets, so a weak session cookie is confirmed honestly, "
            "then hardened on the shield and re-proven."
        ),
    },
    {
        "title": "Class 4 — Privilege Escalation",
        "subtitle": "privilege_escalation — horizontal (BOLA) & vertical",
        "image": "shot_privesc.png", "layout": "right",
        "bullets": [
            ("cyan", "Three-probe differential: control · breach · anon baseline"),
            ("cyan", "CONFIRMED only if control succeeds AND breach granted"),
            ("cyan", "AND the anonymous baseline is denied — no false positives"),
            ("cyan", "A properly owned /orders endpoint → DISPROVED, no finding"),
            ("ok", "Shield denies the cross-boundary request → re-proven"),
        ],
        "notes": (
            "Privilege escalation as a rigorous differential. Sentinel replays the attacker's "
            "session against their own object (control) and a forbidden one (breach), plus an "
            "anonymous baseline; it confirms only when the control succeeds, the breach is "
            "granted, and anonymous is denied. An ownership-checked endpoint yields no finding — "
            "the honest control — then the shield denies the escalation and the judge re-proves."
        ),
    },
    {
        "title": "Class 5 — SQL Injection",
        "subtitle": "injection — a boolean differential, not a signature",
        "image": "shot_injection.png", "layout": "right",
        "bullets": [
            ("cyan", "Benign baseline + length-matched TRUE/FALSE payload pairs"),
            ("cyan", "CONFIRMED only when TRUE ≠ FALSE, anchored to the baseline"),
            ("cyan", "Real SQL toggles the boolean — nothing pattern-matches '1=1'"),
            ("cyan", "A bound/parameterised filter collapses → DISPROVED, no finding"),
            ("ok", "Request-guard shield blocks the payload → re-proven"),
        ],
        "notes": (
            "Injection judged by behaviour, not by regex. Sentinel sends a benign baseline plus "
            "length-matched boolean pairs and confirms only when the TRUE and FALSE arms diverge "
            "while the baseline stays anchored — so the differential emerges from the database "
            "actually evaluating the injected boolean. A parameterised filter collapses the "
            "differential to no finding; the request-guard virtual patch then blocks the payload "
            "and the judge re-proves."
        ),
    },
    {
        "title": "Classes 6–10 — Five More, Same Contract",
        "subtitle": "Every class is its own differential with an explicit anchor",
        "bullets": [
            ("cyan", "SSTI — {{7*7}} evaluates to 49 vs the literal; behaviour, not a regex"),
            ("cyan", "Open redirect — two-probe host differential; off-origin host OBSERVED, never followed"),
            ("cyan", "CORS — two-probe Origin differential; nonce origin reflected AND credentialed"),
            ("cyan", "SSRF — out-of-band callback to Sentinel's OWN loopback collaborator"),
            ("cyan", "Broken auth / JWT — three-probe genuine · forged · absent token differential"),
        ],
        "notes": (
            "Coverage doubled without loosening the bar. SSTI confirms only when injected arithmetic "
            "actually evaluates; open-redirect and CORS each ride a two-probe differential against an "
            "unroutable nonce that is observed or echoed but never contacted; SSRF proves the fetch by "
            "an out-of-band callback to Sentinel's own loopback beacon; broken-auth replays a genuine, "
            "a forged, and an absent token. Ten classes, one contract."
        ),
    },
    {
        "title": "The Hardest Two, Done Honestly",
        "subtitle": "Zero-oracle SSRF · token forgery — and where we stop",
        "bullets": [
            ("cyan", "SSRF is zero-oracle: only Sentinel's own 127.0.0.1 collaborator is ever injected"),
            ("cyan", "— never a metadata IP, an RFC-1918 host, or any third party"),
            ("cyan", "Broken auth: a forged token as SOLE authenticator must be accepted to CONFIRM"),
            ("ok", "Every forgery carries a benign marker — acceptance proves OUR minted token passed"),
            ("warn", "alg-none → full PATCH+PROVE; signed-forgery classes → ADVISORY_ONLY, never a fake fix-proof"),
        ],
        "notes": (
            "The two most sophisticated classes are also where honesty matters most. SSRF needs no "
            "operator oracle — the only URL Sentinel ever injects is its own loopback collaborator, so a "
            "callback is unambiguous and safe. Broken-auth confirms only when a token Sentinel itself "
            "minted (carrying a benign marker) is accepted as the sole authenticator. And when a fix "
            "cannot be proven by a request-guard — HS256-confusion or a cracked weak secret — Sentinel "
            "returns ADVISORY_ONLY with the durable handler-side fix, rather than manufacturing a fix-proof."
        ),
    },
]
SLIDES += [
    {
        "title": "Login Tester — Authenticated Reasoning",
        "subtitle": "Turn anonymous scanning into authenticated proof",
        "image": "shot_login.png", "layout": "right",
        "bullets": [
            ("cyan", "Drives a real browser; you log in (MFA supported)"),
            ("cyan", "Auto-detects completion, captures the live session"),
            ("cyan", "Runs authenticated authz + cookie analysis on it"),
            ("warn", "Credentials via getpass — never stored, never logged"),
            ("warn", "In-memory for the run only, scope-bound to the host"),
        ],
        "notes": (
            "The Login Tester unlocks the authenticated surface. Sentinel opens a visible "
            "browser, waits for you (including MFA), auto-detects success, and reasons over the "
            "real session — with credentials handled in memory only, never persisted or logged."
        ),
    },
    {
        "title": "Chaining — Ingredients, Honestly",
        "subtitle": "We prove the pieces; we don't fabricate the story",
        "bullets": [
            ("ok", "Session captured (real cookies + bearer)"),
            ("ok", "Session cookie proven insecure (insecure_cookie CONFIRMED)"),
            ("ok", "Resource reachable as that principal (authz CONFIRMED)"),
            ("warn", "Presented as co-occurring evidence under one session"),
            ("warn", "Full causal chaining stays the honestly-labeled frontier"),
        ],
        "notes": (
            "Chaining is where scanners over-claim. Sentinel proves each ingredient of a session-"
            "theft→authenticated-access chain independently and presents them together — but it "
            "does not auto-compose a causal exploit narrative. Honesty is the moat."
        ),
    },
    {
        "title": "Target-Agnostic — Two Live Stacks",
        "subtitle": "Zero target-specific code",
        "image": "shot_outcome.png", "layout": "right",
        "bullets": [
            "Proven on OWASP Juice Shop and VAmPI",
            "Same engine, same judges, different oracles",
            "Oracles are operator data, not engine logic",
            "New target = new policy file, no code change",
            "The research frontier is reported, never faked",
        ],
        "notes": (
            "Target-agnosticism is structural: the engine carries no target-specific code, only "
            "the operator's oracle changes. We've run the full loop on two independent stacks — "
            "Juice Shop and VAmPI — to prove it generalizes."
        ),
    },
]
SLIDES += [
    {
        "title": "Test Results & Validation",
        "subtitle": "Green, deterministic, and honest about the one skip",
        "stats": [("255", "tests"), ("254", "passing"), ("10", "vuln classes"),
                  ("2", "live stacks")],
        "bullets": [
            ("ok", "254 passing, 1 skipped (the opt-in live-browser path)"),
            ("ok", "Full offline suite is network-free and runs in seconds"),
            ("ok", "Each class: parse → CONFIRM → PATCH → PROVE, plus isolation"),
            ("ok", "Live E2E validated on Juice Shop and VAmPI over real sockets"),
            ("cyan", "Compliant-control tests assert the honest 'no finding' path"),
        ],
        "notes": (
            "Validation is first-class. The offline suite is deterministic and network-free; the "
            "single skip is the browser test that only runs with the opt-in extra. Tests cover "
            "confirmation, remediation proof, isolation, and — crucially — that compliant "
            "controls produce no finding."
        ),
    },
    {
        "title": "Lightweight & Local",
        "subtitle": "Runs on a laptop; nothing leaves the box",
        "bullets": [
            "Local model via Ollama (qwen3:4b) — no cloud dependency",
            "Core install is lightweight; Playwright is an opt-in extra",
            "No customer data or target traffic egresses anywhere",
            "Deterministic core runs even with the model offline",
            "Deploys air-gapped — ideal for regulated environments",
        ],
        "notes": (
            "Local-first is a feature, not a limitation. The reasoning model runs on-box, the "
            "deterministic core works even if the model is down, and nothing egresses — which "
            "makes Sentinel deployable in air-gapped and regulated settings."
        ),
    },
    {
        "title": "Safety & Bounded Autonomy",
        "subtitle": "Safe by construction",
        "bullets": [
            ("ok", "Pre-connection scope guard — out-of-scope probes refused"),
            ("ok", "Probes are local, bounded, and non-destructive"),
            ("ok", "The enforcer shield is ephemeral and loopback-only (SSRF-safe)"),
            ("ok", "Nothing both acts and judges — stages are separated"),
            ("warn", "Credentials never persisted or logged; in-memory per run"),
        ],
        "notes": (
            "Autonomy is bounded by design: scope is enforced before a socket opens, probes are "
            "non-destructive, the shield is loopback-only, and the acting and judging stages are "
            "strictly separated. Credential handling is memory-only."
        ),
    },
]
SLIDES += [
    {
        "title": "Market Landscape",
        "subtitle": "Where autonomous, provable reasoning fits",
        "cols": ["Signature\nscanners", "DAST /\nfuzzers", "Sentinel"],
        "rows": [
            ("Finds authorization / BOLA flaws", ["no", "partial", "yes"]),
            ("Verdict is a reproduced contradiction", ["no", "no", "yes"]),
            ("Proves the fix, live", ["no", "no", "yes"]),
            ("Autonomous, self-directed research", ["no", "partial", "yes"]),
            ("Local-first, no data egress", ["partial", "partial", "yes"]),
            ("Every decision auditable", ["partial", "no", "yes"]),
        ],
        "notes": (
            "Category-level, not a bake-off: signature scanners are fast but noisy and miss "
            "authorization; DAST and fuzzers explore but don't reason about intent or prove "
            "fixes. Sentinel's column is defined by reproduced verdicts, live fix-proofs, and "
            "full auditability."
        ),
    },
    {
        "title": "Why Sentinel Wins",
        "subtitle": "The defensible core",
        "bullets": [
            ("cyan", "Trust: a finding is a reproduced contradiction, provably"),
            ("cyan", "Closure: it doesn't just find — it patches and re-proves"),
            ("cyan", "Reach: the classes signatures miss (authz, session)"),
            ("cyan", "Honesty: compliant controls and frontiers labeled, not faked"),
            ("cyan", "Deployability: local, auditable, air-gap-ready"),
        ],
        "notes": (
            "The wedge is trust plus closure: verdicts you can believe and fixes you can prove, "
            "in the classes that matter most, deployable where the sensitive targets actually "
            "live. Honesty about limits is itself a differentiator."
        ),
    },
    {
        "title": "Status & Roadmap",
        "subtitle": "Honest today; a clear path to more",
        "bullets": [
            ("ok", "Today: 10 classes live end-to-end on 2 stacks; 254 tests"),
            ("ok", "Today: Login Tester + authenticated reasoning + chaining ingredients"),
            ("next", "Next: provable 2-link chaining (SQLi ⇒ IDOR), decoy-gated"),
            ("next", "Next: the last two OWASP classes — reflected XSS + path traversal"),
            ("next", "Next: CI/CD gate + policy library for common stacks"),
        ],
        "notes": (
            "We're candid about the line between done and next. Ten classes, two targets, the "
            "Login Tester, and chaining ingredients ship today; provable 2-link chaining, the last "
            "two OWASP classes (reflected XSS + path traversal), and a CI gate are the roadmap."
        ),
    },
]
SLIDES += [
    {
        "title": "The Ask",
        "subtitle": "$1M to turn a proven engine into a platform",
        "stats": [("$1M", "seed"), ("18", "months runway"), ("10→12+", "vuln classes")],
        "bullets": [
            ("cyan", "Engineering: provable chaining + the last OWASP classes"),
            ("cyan", "Policy library: ready-made oracles for common stacks"),
            ("cyan", "Integrations: CI/CD gate, ticketing, SIEM export"),
            ("cyan", "Design partners: regulated, air-gap-first customers"),
            ("cyan", "Outcome: from autonomous finder to autonomous fixer"),
        ],
        "notes": (
            "We're raising $1M for ~18 months to expand the class coverage, ship full chaining, "
            "build a policy library and CI integrations, and land design partners in regulated "
            "environments where local-first is a hard requirement."
        ),
    },
    {
        "kind": "divider", "title": "Reason. Prove. Remediate.",
        "subtitle": "Autonomous security you can actually trust",
        "notes": "Close: the future of security tooling isn't more signatures — it's reasoning "
                 "that proves itself. That's Sentinel.",
    },
]
# __SLIDES__

FONT_HEAD = "Consolas"
FONT_BODY = "Segoe UI"
FONT_SYM = "Segoe UI Symbol"  # ✓/✗ live here, not in Segoe UI proper


def _txt(slide, left, top, width, height, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tf


def _run(p, text, size, color, font=FONT_BODY, bold=False):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.name = font
    r.font.bold = bold
    return r


def _badge(slide):
    """AI KAVACH chip, top-right."""
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(11.15), Inches(0.5), Inches(1.7), Inches(0.42))
    shp.fill.background()
    shp.line.color.rgb = C_CYAN
    shp.line.width = Pt(1)
    tf = shp.text_frame
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _run(p, "AI KAVACH", 12, C_CYAN, font=FONT_HEAD, bold=True)


def _footer(slide, n):
    tf = _txt(slide, 0.62, 6.98, 9.0, 0.4)
    p = tf.paragraphs[0]
    _run(p, "SENTINEL", 10, C_CYAN, font=FONT_HEAD, bold=True)
    _run(p, "   ·   find → reason → prove → patch → prove", 10, C_DIM, font=FONT_HEAD)
    tfn = _txt(slide, 11.8, 6.98, 1.0, 0.4)
    pn = tfn.paragraphs[0]
    pn.alignment = PP_ALIGN.RIGHT
    _run(pn, f"{n:02d} / {len(SLIDES):02d}", 10, C_DIM, font=FONT_HEAD)
# __BUILD__

_GLYPH = {"ok": ("»", C_GREEN, C_WHITE), "next": ("»", C_AMBER, C_WHITE),
          "cyan": ("»", C_CYAN, C_BODY), "warn": ("»", C_MAGENTA, C_BODY)}


def _content_header(slide, s, mark, rule):
    slide.shapes.add_picture(mark, Inches(0.62), Inches(0.5), width=Inches(0.82))
    _badge(slide)
    tt = _txt(slide, 1.65, 0.52, 9.3, 1.15, anchor=MSO_ANCHOR.MIDDLE)
    _run(tt.paragraphs[0], s["title"], 31, C_WHITE, font=FONT_HEAD, bold=True)
    slide.shapes.add_picture(rule, Inches(0.7), Inches(1.82),
                             width=Inches(3.4), height=Inches(0.062))
    st = _txt(slide, 0.72, 2.02, 11.9, 0.6)
    _run(st.paragraphs[0], s["subtitle"], 16, C_CYAN, font=FONT_HEAD)


def _bullets(slide, bullets, left, top, width, size=18, gap=13):
    bt = _txt(slide, left, top, width, 4.2)
    for i, b in enumerate(bullets):
        kind, text = ("cyan", b) if isinstance(b, str) else b
        glyph, gcol, tcol = _GLYPH[kind]
        p = bt.paragraphs[0] if i == 0 else bt.add_paragraph()
        p.space_after = Pt(gap)
        _run(p, glyph + "  ", size, gcol, font=FONT_BODY, bold=True)
        _run(p, text, size, tcol, font=FONT_BODY)


def _picture(slide, name, box):
    """Fit image into box=(left,top,maxw,maxh) inches, return (left,top,w,h)."""
    left, top, maxw, maxh = box
    im = Image.open(os.path.join(BRAND, name))
    ar = im.width / im.height
    w, h = maxw, maxw / ar
    if h > maxh:
        h, w = maxh, maxh * ar
    x = left + (maxw - w) / 2
    y = top + (maxh - h) / 2
    slide.shapes.add_picture(os.path.join(BRAND, name), Inches(x), Inches(y),
                             width=Inches(w), height=Inches(h))
    return x, y, w, h
# __BUILD2__

def _stats(slide, stats, top):
    n = len(stats)
    gapw = 0.35
    total = 12.0
    cw = (total - gapw * (n - 1)) / n
    for i, (val, lab) in enumerate(stats):
        x = 0.66 + i * (cw + gapw)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(x), Inches(top), Inches(cw), Inches(1.25))
        card.fill.background()
        card.line.color.rgb = C_CYAN
        card.line.width = Pt(1)
        tf = card.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        pv = tf.paragraphs[0]
        pv.alignment = PP_ALIGN.CENTER
        _run(pv, val, 32, C_CYAN, font=FONT_HEAD, bold=True)
        pl = tf.add_paragraph()
        pl.alignment = PP_ALIGN.CENTER
        _run(pl, lab.upper(), 11, C_DIM, font=FONT_HEAD)


_CELL = {"yes": ("✓", C_GREEN), "partial": ("~", C_AMBER), "no": ("✗", C_RED)}


def _compare(slide, cols, rows, top):
    x_lab, w_lab = 0.72, 6.2
    x0, cw = 7.15, 1.85
    # column headers
    for j, c in enumerate(cols):
        hb = _txt(slide, x0 + j * cw, top, cw, 0.7)
        hb.vertical_anchor = MSO_ANCHOR.BOTTOM
        for k, line in enumerate(c.split("\n")):
            p = hb.paragraphs[0] if k == 0 else hb.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            _run(p, line, 12, C_WHITE, font=FONT_HEAD, bold=True)
    y = top + 0.78
    rh = 0.52
    for ri, (label, states) in enumerate(rows):
        if ri % 2 == 0:
            band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(y),
                                          Inches(12.1), Inches(rh))
            band.fill.solid()
            band.fill.fore_color.rgb = RGBColor(0x10, 0x17, 0x24)
            band.line.fill.background()
            band.shadow.inherit = False
        lb = _txt(slide, x_lab, y, w_lab, rh, anchor=MSO_ANCHOR.MIDDLE)
        _run(lb.paragraphs[0], label, 14, C_BODY, font=FONT_BODY)
        for j, stt in enumerate(states):
            glyph, col = _CELL[stt]
            cb = _txt(slide, x0 + j * cw, y, cw, rh, anchor=MSO_ANCHOR.MIDDLE)
            pc = cb.paragraphs[0]
            pc.alignment = PP_ALIGN.CENTER
            _run(pc, glyph, 18, col, font=FONT_SYM, bold=True)
        y += rh
# __BUILD3__

def _divider(slide, s, rule):
    tt = _txt(slide, 1.0, 2.75, 11.333, 1.5, anchor=MSO_ANCHOR.MIDDLE)
    pt = tt.paragraphs[0]
    pt.alignment = PP_ALIGN.CENTER
    _run(pt, s["title"], 40, C_WHITE, font=FONT_HEAD, bold=True)
    slide.shapes.add_picture(rule, Inches((13.333 - 3.4) / 2), Inches(4.35),
                             width=Inches(3.4), height=Inches(0.062))
    st = _txt(slide, 1.0, 4.55, 11.333, 0.8)
    ps = st.paragraphs[0]
    ps.alignment = PP_ALIGN.CENTER
    _run(ps, s["subtitle"], 18, C_CYAN, font=FONT_HEAD)


def build():
    content_bg = os.path.join(BRAND, "deck_bg.png")
    title_bg = os.path.join(BRAND, "deck_bg_title.png")
    divider_bg = os.path.join(BRAND, "deck_bg_divider.png")
    rule = os.path.join(BRAND, "deck_rule.png")
    mark = os.path.join(BRAND, "sentinel_mark.png")
    render_content_bg(content_bg)
    render_title_bg(title_bg)
    render_divider_bg(divider_bg)
    render_rule(rule)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for idx, s in enumerate(SLIDES, start=1):
        kind = s.get("kind", "content")
        slide = prs.slides.add_slide(blank)
        bg = {"title": title_bg, "divider": divider_bg}.get(kind, content_bg)
        slide.shapes.add_picture(bg, 0, 0, width=prs.slide_width, height=prs.slide_height)

        if kind == "title":
            lock = os.path.join(BRAND, "sentinel_logo_transparent.png")
            lw = 9.2
            im = Image.open(lock)
            lh = lw * im.height / im.width
            slide.shapes.add_picture(lock, Inches((13.333 - lw) / 2), Inches(1.55),
                                     width=Inches(lw))
            sub = _txt(slide, 1.5, 1.7 + lh, 10.333, 0.6)
            ps = sub.paragraphs[0]
            ps.alignment = PP_ALIGN.CENTER
            _run(ps, s["subtitle"], 18, C_CYAN, font=FONT_HEAD)
            body = _txt(slide, 1.9, 2.35 + lh, 9.533, 2.2)
            for i, b in enumerate(s["bullets"]):
                p = body.paragraphs[0] if i == 0 else body.add_paragraph()
                p.alignment = PP_ALIGN.CENTER
                p.space_after = Pt(7)
                _run(p, b, 15, C_BODY, font=FONT_BODY)
            foot = _txt(slide, 0, 6.85, 13.333, 0.5)
            pf = foot.paragraphs[0]
            pf.alignment = PP_ALIGN.CENTER
            _run(pf, "AI KAVACH CHALLENGE  ·  AUTONOMOUS CYBER-REASONING", 11,
                 C_DIM, font=FONT_HEAD, bold=True)
        elif kind == "divider":
            _divider(slide, s, rule)
        else:
            _render_content(slide, s, mark, rule, idx)

        slide.notes_slide.notes_text_frame.text = s["notes"]

    out = os.path.join(DECK, "Sentinel_AI_Kavach.pptx")
    prs.save(out)
    print("wrote", out)
    print("slides:", len(prs.slides._sldIdLst))
# __BUILD4__

def _render_content(slide, s, mark, rule, idx):
    _content_header(slide, s, mark, rule)
    image = s.get("image")
    layout = s.get("layout", "right")

    if "cols" in s:  # market matrix
        _compare(slide, s["cols"], s["rows"], top=2.62)
    elif "stats" in s:
        _stats(slide, s["stats"], top=2.62)
        if s.get("bullets"):
            _bullets(slide, s["bullets"], 0.78, 4.2, 11.8, size=17, gap=10)
    elif image and layout == "full":
        _, iy, _, ih = _picture(slide, image, (1.0, 2.55, 11.3, 4.05))
        cap = s.get("caption")
        if cap:
            cb = _txt(slide, 1.0, iy + ih + 0.06, 11.3, 0.4)
            pc = cb.paragraphs[0]
            pc.alignment = PP_ALIGN.CENTER
            _run(pc, cap, 12, C_DIM, font=FONT_HEAD)
    elif image:  # right layout: bullets left, image right
        _picture(slide, image, (7.0, 2.55, 5.85, 4.15))
        if s.get("bullets"):
            _bullets(slide, s["bullets"], 0.78, 2.72, 5.95, size=16, gap=12)
    else:
        _bullets(slide, s["bullets"], 0.78, 2.95, 11.8, size=18, gap=15)

    _footer(slide, idx)


if __name__ == "__main__":
    build()















